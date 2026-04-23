from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Color Scheme - Professional Light Theme
PRIMARY_COLOR = RGBColor(0, 102, 204)      # Professional Blue
SECONDARY_COLOR = RGBColor(51, 153, 204)   # Light Blue
ACCENT_COLOR = RGBColor(255, 102, 0)       # Orange Accent
TEXT_DARK = RGBColor(40, 40, 40)           # Dark Gray
TEXT_LIGHT = RGBColor(100, 100, 100)       # Medium Gray
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(245, 248, 250)         # Very Light Blue

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = "AUTOHAND"
    title_p.font.size = Pt(88)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "AI-Based Healthcare Automation Agent"
    subtitle_p.font.size = Pt(36)
    subtitle_p.font.color.rgb = WHITE
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    footer_frame = footer_box.text_frame
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "Team Members, College"
    footer_p.font.size = Pt(20)
    footer_p.font.color.rgb = RGBColor(200, 220, 240)
    footer_p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Use Case Problem
    slide2 = add_content_slide(prs, "Use Case: Current Challenges")
    add_bullet_points(slide2, [
        "Manual patient data entry",
        "Delay in generating reports",
        "No real-time response capability"
    ], Inches(1), Inches(2))
    
    # Add visual - Doctor with arrows
    add_flowchart_visual(slide2)
    
    # Slide 3: Use Case Explanation
    slide3 = add_content_slide(prs, "Problems We're Solving")
    add_bullet_points(slide3, [
        "Manual Excel updates - Time-consuming",
        "Email reports individually - Error-prone",
        "No automation - Operational bottleneck"
    ], Inches(1), Inches(2))
    
    # Add problem icon
    add_visual_elements(slide3, "⚠", Inches(8), Inches(3), 60)
    
    # Slide 4: Philips Overview
    slide4 = add_content_slide(prs, "Healthcare Innovation Partner")
    add_bullet_points(slide4, [
        "Global healthcare technology leader",
        "Focus on AI & digital transformation",
        "Goal: Improve efficiency and patient outcomes"
    ], Inches(1), Inches(2))
    
    # Slide 5: Mind Map - Problem to Solution
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide5)
    add_title(slide5, "AUTOHAND: Solution Overview")
    add_mind_map(slide5, "AUTOHAND", [
        "Data Handling",
        "Automation",
        "AI Analysis",
        "Communication",
        "Emergency Response"
    ])
    
    # Slide 6: How AUTOHAND Helps
    slide6 = add_content_slide(prs, "AUTOHAND Benefits")
    benefits = [
        "↓ Reduces manual effort by 80%",
        "⏱ Saves time - Automates routine tasks",
        "📊 Automates report generation",
        "🚨 Sends instant alerts to healthcare team",
        "⚡ Detects emergencies in real-time"
    ]
    add_bullet_points(slide6, benefits, Inches(1.5), Inches(2), Pt(22))
    
    # Slide 7: Operations Mind Map
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide7)
    add_title(slide7, "Agent Operations")
    add_mind_map(slide7, "Agent\nOperations", [
        "App Control",
        "Excel Data",
        "Email Alerts",
        "AI Analysis",
        "System Tasks"
    ])
    
    # Slide 8: Architecture
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide8)
    add_title(slide8, "System Architecture")
    add_architecture_flow(slide8)
    
    # Slide 9: Tech Stack
    slide9 = add_content_slide(prs, "Technology Stack")
    tech_stack = [
        "Backend: Python 3.12 + Flask",
        "Data Processing: Pandas, OpenPyXL",
        "System Integration: PyWin32",
        "AI Engine: Google Gemini API",
        "Frontend: HTML5, CSS3, JavaScript"
    ]
    add_bullet_points(slide9, tech_stack, Inches(1.5), Inches(2), Pt(22))
    
    # Slide 10: Workflow
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide10)
    add_title(slide10, "Execution Workflow")
    add_workflow_steps(slide10)
    
    # Slide 11: Demo Preview
    slide11 = add_content_slide(prs, "Live Demo Workflow")
    demo_steps = [
        "1. User input healthcare command",
        "2. AI processes natural language",
        "3. Excel file updated automatically",
        "4. Data analysis performed",
        "5. Email alerts sent to team"
    ]
    add_bullet_points(slide11, demo_steps, Inches(1.5), Inches(2), Pt(22))
    
    # Slide 12: Benefits Summary
    slide12 = add_content_slide(prs, "Impact & Benefits")
    impact = [
        "💰 Cost Reduction: 70% less manual work",
        "✅ Quality: 99% accuracy in automation",
        "⚡ Speed: Real-time response capability",
        "📈 Scalability: Handle thousands of records"
    ]
    add_bullet_points(slide12, impact, Inches(1.5), Inches(2), Pt(22))
    
    # Slide 13: Future Roadmap
    slide13 = add_content_slide(prs, "Future Enhancements")
    future = [
        "🎙 Voice-controlled commands",
        "📊 Interactive dashboard & analytics",
        "🔔 Real-time monitoring system",
        "🔐 Enhanced security & compliance"
    ]
    add_bullet_points(slide13, future, Inches(1.5), Inches(2), Pt(22))
    
    # Slide 14: Conclusion
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide14.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    
    conclusion_box = slide14.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    conclusion_frame = conclusion_box.text_frame
    conclusion_frame.word_wrap = True
    
    p = conclusion_frame.paragraphs[0]
    p.text = "AUTOHAND: Transforming Healthcare Automation"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = conclusion_frame.add_paragraph()
    p2.text = "\nAI-Powered Intelligence for Modern Healthcare"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(200, 220, 240)
    p2.alignment = PP_ALIGN.CENTER
    
    # Slide 15: Thank You
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide15.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    thankyou_box = slide15.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    thankyou_frame = thankyou_box.text_frame
    thankyou_frame.word_wrap = True
    
    p = thankyou_frame.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = thankyou_frame.add_paragraph()
    p2.text = "\nQuestions & Discussion"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(200, 220, 240)
    p2.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = "AUTOHAND_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created successfully: {output_path}")
    return output_path

def add_slide_background(slide):
    """Add light background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG

def add_title(slide, title_text):
    """Add title to slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR

def add_content_slide(prs, title_text):
    """Create a standard content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)
    add_title(slide, title_text)
    return slide

def add_bullet_points(slide, points, left, top, font_size=Pt(24)):
    """Add bullet points to slide"""
    textbox = slide.shapes.add_textbox(left, top, Inches(8), Inches(4.5))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = font_size
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0

def add_mind_map(slide, center, branches):
    """Add mind map with center and branches"""
    center_x = Inches(5)
    center_y = Inches(3.5)
    
    # Center circle
    center_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        center_x - Inches(0.7), center_y - Inches(0.5),
        Inches(1.4), Inches(1)
    )
    center_circle.fill.solid()
    center_circle.fill.fore_color.rgb = PRIMARY_COLOR
    center_circle.line.color.rgb = PRIMARY_COLOR
    
    center_text = center_circle.text_frame
    center_text.text = center
    center_text.word_wrap = True
    center_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    for para in center_text.paragraphs:
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    
    # Branches
    num_branches = len(branches)
    angle_step = 360 / num_branches
    
    for i, branch in enumerate(branches):
        angle = i * angle_step
        
        # Calculate position
        radius = Inches(2.2)
        import math
        rad = math.radians(angle)
        x = center_x + radius * math.cos(rad) - Inches(0.5)
        y = center_y + radius * math.sin(rad) - Inches(0.35)
        
        # Branch circle
        branch_circle = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, Inches(1), Inches(0.7)
        )
        branch_circle.fill.solid()
        branch_circle.fill.fore_color.rgb = SECONDARY_COLOR
        branch_circle.line.color.rgb = PRIMARY_COLOR
        branch_circle.line.width = Pt(2)
        
        # Branch text
        branch_text = branch_circle.text_frame
        branch_text.text = branch
        branch_text.word_wrap = True
        branch_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for para in branch_text.paragraphs:
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
        
        # Connector line
        connector = slide.shapes.add_connector(1, int(center_x), int(center_y), int(x + Inches(0.5)), int(y + Inches(0.35)))
        connector.line.color.rgb = PRIMARY_COLOR
        connector.line.width = Pt(2)

def add_flowchart_visual(slide):
    """Add flowchart visual for problem flow"""
    boxes = ["Doctor", "Manual\nWork", "Delay", "Risk"]
    colors = [PRIMARY_COLOR, ACCENT_COLOR, RGBColor(255, 102, 102), RGBColor(204, 0, 0)]
    
    start_x = Inches(1.5)
    for i, (box_text, color) in enumerate(zip(boxes, colors)):
        x = start_x + i * Inches(2)
        y = Inches(4.5)
        
        # Box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, Inches(1.5), Inches(1)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(100, 100, 100)
        
        # Text
        text_frame = shape.text_frame
        text_frame.text = box_text
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for para in text_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
        
        # Arrow
        if i < len(boxes) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + Inches(1.6), y + Inches(0.4),
                Inches(0.5), Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEXT_LIGHT
            arrow.line.color.rgb = TEXT_LIGHT

def add_visual_elements(slide, symbol, x, y, size):
    """Add visual symbols to slide"""
    textbox = slide.shapes.add_textbox(x, y, Inches(1.5), Inches(1))
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = symbol
    p.font.size = Pt(size)
    p.alignment = PP_ALIGN.CENTER

def add_architecture_flow(slide):
    """Add architecture diagram"""
    components = ["User", "AI Agent", "Decision\nEngine", "Modules", "Output"]
    
    start_x = Inches(0.8)
    for i, comp in enumerate(components):
        x = start_x + i * Inches(1.8)
        y = Inches(3.5)
        
        # Component box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, Inches(1.5), Inches(1)
        )
        shape.fill.solid()
        
        if i == 0 or i == len(components) - 1:
            shape.fill.fore_color.rgb = ACCENT_COLOR
        else:
            shape.fill.fore_color.rgb = PRIMARY_COLOR
        
        shape.line.color.rgb = PRIMARY_COLOR
        shape.line.width = Pt(2)
        
        # Text
        text_frame = shape.text_frame
        text_frame.text = comp
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for para in text_frame.paragraphs:
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
        
        # Arrow
        if i < len(components) - 1:
            arrow = slide.shapes.add_connector(1, int(x + Inches(1.5)), int(y + Inches(0.5)), 
                                              int(x + Inches(1.8)), int(y + Inches(0.5)))
            arrow.line.color.rgb = PRIMARY_COLOR
            arrow.line.width = Pt(3)

def add_workflow_steps(slide):
    """Add workflow steps"""
    steps = ["Command", "AI\nProcessing", "Decision", "Action", "Output"]
    colors = [PRIMARY_COLOR, SECONDARY_COLOR, ACCENT_COLOR, RGBColor(100, 180, 100), RGBColor(150, 100, 200)]
    
    start_x = Inches(0.8)
    for i, (step, color) in enumerate(zip(steps, colors)):
        x = start_x + i * Inches(1.8)
        y = Inches(3)
        
        # Step circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x, y, Inches(1.4), Inches(1.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = RGBColor(100, 100, 100)
        circle.line.width = Pt(2)
        
        # Text
        text_frame = circle.text_frame
        text_frame.text = step
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for para in text_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
        
        # Arrow
        if i < len(steps) - 1:
            arrow = slide.shapes.add_connector(1, int(x + Inches(1.4)), int(y + Inches(0.7)), 
                                              int(x + Inches(1.8)), int(y + Inches(0.7)))
            arrow.line.color.rgb = RGBColor(150, 150, 150)
            arrow.line.width = Pt(3)

if __name__ == "__main__":
    create_presentation()
